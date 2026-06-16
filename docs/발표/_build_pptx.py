# -*- coding: utf-8 -*-
"""
한화 ATLAS 발표 PPTX 빌더 (제목 + 문제제기 → 데모 전환까지)
- 바이브코딩 경진대회(주식운용 아이디어) 5분 발표용
- 한화 팔레트(오렌지/브라운/베이지) · 다크 프리미엄 톤 · Malgun Gothic
실행: C:/Users/infomax/AppData/Local/Programs/Python/Python313/python.exe _build_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---- 한화 브랜드 팔레트 ----
ORANGE      = RGBColor(0xF3, 0x73, 0x21)   # 한화 메인 오렌지
ORANGE_SOFT = RGBColor(0xFF, 0x9A, 0x52)   # 밝은 오렌지(보조)
BG_DARK     = RGBColor(0x2B, 0x1D, 0x16)   # 따뜻한 다크 브라운 배경
BG_DARKER   = RGBColor(0x1C, 0x12, 0x0D)   # 더 어두운 브라운
BROWN_MID   = RGBColor(0x4A, 0x33, 0x26)   # 중간 브라운(카드)
BEIGE       = RGBColor(0xF3, 0xE7, 0xDA)   # 밝은 베이지 텍스트
BEIGE_DIM   = RGBColor(0xC9, 0xB6, 0xA4)   # 흐린 베이지(서브텍스트)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Malgun Gothic"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    """슬라이드 전체 배경 사각형."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    # 맨 뒤로
    sp = shp._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return shp


def rect(slide, x, y, w, h, color, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def bar(slide, x, y, w, h, color):
    """장식용 색 막대(테두리 없음)."""
    return rect(slide, x, y, w, h, color)


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=Pt(6), line_spacing=1.05):
    """runs: list of paragraphs, each paragraph is list of (text, size, color, bold, font)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, size, color, bold, *rest) in para:
            fn = rest[0] if rest else FONT
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = fn
            # 한글 폰트 강제 지정(eastasian)
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {})
                rPr.append(ea)
            ea.set('typeface', fn)
    return tb


def page_num(slide, n, total):
    txt(slide, SW - Inches(1.2), SH - Inches(0.55), Inches(0.9), Inches(0.4),
        [[(f"{n} / {total}", 11, BEIGE_DIM, False)]], align=PP_ALIGN.RIGHT)


def brand_tag(slide):
    txt(slide, Inches(0.6), SH - Inches(0.55), Inches(4.0), Inches(0.4),
        [[("한화 ATLAS", 11, ORANGE, True), ("  ·  AI Investment Desk OS", 11, BEIGE_DIM, False)]])


TOTAL = 6

# ======================================================================
# Slide 1 — 타이틀
# ======================================================================
s = add_slide()
fill_bg(s, BG_DARK)
# 좌측 오렌지 액센트 기둥
bar(s, 0, 0, Inches(0.28), SH, ORANGE)
# 상단 라벨
txt(s, Inches(1.0), Inches(1.15), Inches(11), Inches(0.5),
    [[("한화손해보험 운용본부  ·  AI 주식 운용 데스크 OS", 16, ORANGE_SOFT, True)]])
# 메인 타이틀
txt(s, Inches(0.95), Inches(1.95), Inches(11.4), Inches(2.2),
    [[("한화 ", 88, BEIGE, True), ("ATLAS", 88, ORANGE, True)]],
    anchor=MSO_ANCHOR.TOP)
# 한 줄 소개
txt(s, Inches(1.0), Inches(3.85), Inches(11.2), Inches(1.2),
    [[("운용역의 하루를 함께 도는 ", 26, BEIGE, False),
      ("AI 멀티에이전트 운용 데스크", 26, WHITE, True)],
     [("시황 브리핑 · AI 투자위원회 · 아이디어랩 · 시장 대시보드를 하나의 OS로", 17, BEIGE_DIM, False)]],
    line_spacing=1.2)
# 라이브 URL 박스
box = rect(s, Inches(1.0), Inches(5.55), Inches(7.4), Inches(0.72), BG_DARKER)
box.line.color.rgb = ORANGE
box.line.width = Pt(1.25)
txt(s, Inches(1.25), Inches(5.62), Inches(7.0), Inches(0.6),
    [[("LIVE  ", 15, ORANGE, True),
      ("hanwha-atlas-production.up.railway.app", 15, BEIGE, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
# 발표자 자리
txt(s, Inches(1.0), Inches(6.55), Inches(11), Inches(0.5),
    [[("바이브코딩 경진대회 (주식운용 아이디어)   |   발표자: ____________   |   2026", 14, BEIGE_DIM, False)]])
brand_tag(s)
page_num(s, 1, TOTAL)

# ======================================================================
# Slide 2 — 문제제기 (1) 정보 과부하
# ======================================================================
s = add_slide()
fill_bg(s, BG_DARKER)
bar(s, 0, 0, SW, Inches(0.16), ORANGE)
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.55),
    [[("PROBLEM", 15, ORANGE, True), ("   운용역의 하루는 ‘정보 전쟁’입니다", 15, BEIGE_DIM, True)]])
txt(s, Inches(0.85), Inches(1.35), Inches(11.6), Inches(1.7),
    [[("쏟아지는 정보, ", 46, BEIGE, True), ("소화할 시간은 없다", 46, ORANGE, True)]],
    line_spacing=1.05)
# 3개 페인 카드
cards = [
    ("간밤 미국장", "S&P·나스닥·반도체·환율·금리를 일일이 확인", ""),
    ("매크로 이벤트", "CPI·FOMC·실적·지정학…쉴 새 없는 헤드라인", ""),
    ("국내장 개장", "장 열리기 전에 ‘오늘의 그림’을 다 그려야 한다", ""),
]
cx = Inches(0.85)
cw = Inches(3.75)
gap = Inches(0.13)
cy = Inches(3.35)
ch = Inches(2.35)
for i, (h, d, _) in enumerate(cards):
    x = Emu(int(cx) + i * (int(cw) + int(gap)))
    c = rect(s, x, cy, cw, ch, BROWN_MID)
    c.line.color.rgb = ORANGE
    c.line.width = Pt(0.75)
    bar(s, x, cy, cw, Inches(0.10), ORANGE)
    txt(s, Emu(int(x) + Inches(0.28)), Emu(int(cy) + Inches(0.4)),
        Emu(int(cw) - Inches(0.5)), Inches(1.7),
        [[(h, 23, ORANGE_SOFT, True)],
         [("", 6, BEIGE, False)],
         [(d, 16, BEIGE, False)]], line_spacing=1.2)
txt(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.6),
    [[("→ 정보는 넘치는데, ", 18, BEIGE_DIM, False),
      ("판단으로 바꿀 시간이 부족하다", 18, WHITE, True)]])
brand_tag(s)
page_num(s, 2, TOTAL)

# ======================================================================
# Slide 3 — 문제제기 (2) 반복되는 수작업 시황
# ======================================================================
s = add_slide()
fill_bg(s, BG_DARK)
bar(s, 0, 0, SW, Inches(0.16), ORANGE)
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.55),
    [[("PROBLEM", 15, ORANGE, True), ("   매일, 하루 세 번, 손으로 다시 쓴다", 15, BEIGE_DIM, True)]])
txt(s, Inches(0.85), Inches(1.35), Inches(11.6), Inches(1.7),
    [[("장전·장중·마감마다 ", 44, BEIGE, True), ("반복되는 수작업", 44, ORANGE, True)]],
    line_spacing=1.05)
# 타임라인 3슬롯
slots = [
    ("장전", "07–08시", "간밤 글로벌 → 오늘 전략을 직접 정리"),
    ("장중", "장중", "급변 이벤트마다 메모를 다시 손질"),
    ("마감", "15:30+", "하루를 복기하고 리포트로 또 정리"),
]
ty = Inches(3.3)
tw = Inches(3.8)
tgap = Inches(0.16)
tx0 = Inches(0.85)
for i, (label, when, desc) in enumerate(slots):
    x = Emu(int(tx0) + i * (int(tw) + int(tgap)))
    c = rect(s, x, ty, tw, Inches(2.4), BG_DARKER)
    c.line.color.rgb = BROWN_MID
    c.line.width = Pt(1)
    # 원형 슬롯 배지
    badge = slide_badge = rect(s, Emu(int(x) + Inches(0.28)), Emu(int(ty) + Inches(0.28)),
                               Inches(1.15), Inches(0.6), ORANGE)
    txt(s, Emu(int(x) + Inches(0.28)), Emu(int(ty) + Inches(0.30)),
        Inches(1.15), Inches(0.55),
        [[(label, 18, BG_DARKER, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Emu(int(x) + Inches(1.6)), Emu(int(ty) + Inches(0.34)),
        Inches(2.0), Inches(0.5), [[(when, 15, ORANGE_SOFT, True)]])
    txt(s, Emu(int(x) + Inches(0.3)), Emu(int(ty) + Inches(1.15)),
        Emu(int(tw) - Inches(0.55)), Inches(1.1),
        [[(desc, 16, BEIGE, False)]], line_spacing=1.2)
txt(s, Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.6),
    [[("→ 가치 있는 건 ", 18, BEIGE_DIM, False),
      ("판단", 18, ORANGE, True),
      ("인데, 시간은 ", 18, BEIGE_DIM, False),
      ("정리", 18, WHITE, True),
      ("에 쓰인다", 18, BEIGE_DIM, False)]])
brand_tag(s)
page_num(s, 3, TOTAL)

# ======================================================================
# Slide 4 — 문제제기 (3) 의사결정 근거 & 섹터 발굴
# ======================================================================
s = add_slide()
fill_bg(s, BG_DARKER)
bar(s, 0, 0, SW, Inches(0.16), ORANGE)
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.55),
    [[("PROBLEM", 15, ORANGE, True), ("   ‘왜 샀나’를 매번 새로 설명해야 한다", 15, BEIGE_DIM, True)]])
txt(s, Inches(0.85), Inches(1.35), Inches(11.6), Inches(1.7),
    [[("근거는 흩어지고, ", 44, BEIGE, True), ("발굴은 막막하다", 44, ORANGE, True)]],
    line_spacing=1.05)
# 2개 큰 페인 카드
big = [
    ("종목 의사결정의 근거 부족·일관성 결여",
     ["기술·재무·뉴스·심리가 머릿속에만 흩어져 있다",
      "사람마다, 날마다 판단 기준이 달라진다",
      "‘왜 BUY/SELL인가’를 기록·재현하기 어렵다"]),
    ("섹터 발굴의 막막함",
     ["“반도체 말고 지금 뭐가 좋지?”에 답이 없다",
      "뉴스플로우·거시·수급을 한 번에 못 본다",
      "아이디어는 감(感)에 의존하고 휘발된다"]),
]
by = Inches(3.25)
bw = Inches(5.75)
bgap = Inches(0.3)
bx0 = Inches(0.85)
for i, (h, items) in enumerate(big):
    x = Emu(int(bx0) + i * (int(bw) + int(bgap)))
    c = rect(s, x, by, bw, Inches(2.85), BROWN_MID)
    c.line.color.rgb = ORANGE
    c.line.width = Pt(0.75)
    bar(s, x, by, Inches(0.10), Inches(2.85), ORANGE)
    txt(s, Emu(int(x) + Inches(0.4)), Emu(int(by) + Inches(0.32)),
        Emu(int(bw) - Inches(0.7)), Inches(0.9),
        [[(h, 21, ORANGE_SOFT, True)]], line_spacing=1.1)
    paras = [[("•  ", 16, ORANGE, True), (it, 16, BEIGE, False)] for it in items]
    txt(s, Emu(int(x) + Inches(0.42)), Emu(int(by) + Inches(1.3)),
        Emu(int(bw) - Inches(0.7)), Inches(1.5), paras, line_spacing=1.25, space_after=Pt(8))
txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
    [[("→ 결국 필요한 건 ", 18, BEIGE_DIM, False),
      ("정리·토론·발굴을 대신 돌려주는 ‘운용 데스크’", 18, WHITE, True)]])
brand_tag(s)
page_num(s, 4, TOTAL)

# ======================================================================
# Slide 5 — 해법 한 줄 + 4기능 요약(데모 전환 직전 다리)
# ======================================================================
s = add_slide()
fill_bg(s, BG_DARK)
bar(s, 0, 0, Inches(0.28), SH, ORANGE)
txt(s, Inches(0.95), Inches(0.7), Inches(11.5), Inches(0.55),
    [[("SOLUTION", 15, ORANGE, True), ("   그래서 만들었습니다", 15, BEIGE_DIM, True)]])
txt(s, Inches(0.9), Inches(1.3), Inches(11.6), Inches(1.6),
    [[("한 화면에서 도는 ", 42, BEIGE, True), ("AI 운용 데스크 OS", 42, ORANGE, True)]],
    line_spacing=1.05)
feats = [
    ("01  시황 브리핑 에이전트", "장전·장중·마감 3슬롯. 미국장→매크로→국내장 인과 내러티브 + 업종 로테이션 + KOSPI grounding. 전문가급 PNG 리포트 렌더 + 텔레그램 자동 발송 · 24h 캐싱."),
    ("02  AI 투자위원회", "종목 입력 → 14개 AI 에이전트(기술·심리·뉴스·재무 → Bull/Bear 토론 → 3-way 리스크 심의 → 트레이더·의장 최종결정). 라이브피드 토론 + 9개 리포트 + BUY/SELL/HOLD."),
    ("03  AI 아이디어랩", "키워드(예: ‘반도체 제외 유망섹터’) → 뉴스플로우·거시·마켓무브로 지금 적합한 섹터를 동적 발굴하는 멀티에이전트 발굴위원회. 후보 종목 + 회의록."),
    ("04  시장 대시보드", "KOSPI 전종목 실시간 시세 · GICS 11개 대분류 시장 히트맵 · 섹터 등락 · 손익(mock)."),
]
fy = Inches(2.85)
fw = Inches(5.75)
fh = Inches(1.85)
fgx = Inches(0.3)
fgy = Inches(0.2)
fx0 = Inches(0.9)
for i, (h, d) in enumerate(feats):
    col = i % 2
    row = i // 2
    x = Emu(int(fx0) + col * (int(fw) + int(fgx)))
    y = Emu(int(fy) + row * (int(fh) + int(fgy)))
    c = rect(s, x, y, fw, fh, BG_DARKER)
    c.line.color.rgb = BROWN_MID
    c.line.width = Pt(1)
    bar(s, x, y, Inches(0.09), fh, ORANGE)
    txt(s, Emu(int(x) + Inches(0.35)), Emu(int(y) + Inches(0.22)),
        Emu(int(fw) - Inches(0.6)), Inches(0.5),
        [[(h, 19, ORANGE_SOFT, True)]])
    txt(s, Emu(int(x) + Inches(0.35)), Emu(int(y) + Inches(0.78)),
        Emu(int(fw) - Inches(0.6)), Inches(1.0),
        [[(d, 13.5, BEIGE, False)]], line_spacing=1.18)
txt(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
    [[("LLM: MiMo 추론모델  ·  데이터: 네이버/yfinance 무키 소스  ·  결과 영구볼륨 24h 영속  ·  Railway 단일 서비스(백+프론트)", 12.5, BEIGE_DIM, False)]])
brand_tag(s)
page_num(s, 5, TOTAL)

# ======================================================================
# Slide 6 — 데모 전환
# ======================================================================
s = add_slide()
fill_bg(s, BG_DARKER)
# 중앙 오렌지 굵은 가로 라인
bar(s, Inches(0.0), Inches(2.55), SW, Inches(0.06), ORANGE)
txt(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(0.6),
    [[("NOW, LIVE", 16, ORANGE, True)]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.7), Inches(2.75), Inches(12), Inches(1.8),
    [[("백 마디 설명보다,", 50, BEIGE, True)],
     [("지금 ", 50, BEIGE, True), ("직접 돌려서 보여드립니다", 50, ORANGE, True)]],
    align=PP_ALIGN.CENTER, line_spacing=1.1)
# 데모 순서 칩
chips = ["시장 대시보드 히트맵", "시황 에이전트 · 장전 생성 → PNG", "AI 투자위원회 · 라이브 토론", "아이디어랩 · 섹터 발굴"]
chip_y = Inches(5.05)
total_w = Inches(11.8)
n = len(chips)
chip_w = Inches(2.78)
chip_gap = Emu((int(total_w) - n * int(chip_w)) // (n - 1))
start_x = Emu((int(SW) - int(total_w)) // 2)
for i, ctext in enumerate(chips):
    x = Emu(int(start_x) + i * (int(chip_w) + int(chip_gap)))
    c = rect(s, x, chip_y, chip_w, Inches(0.95), BROWN_MID)
    c.line.color.rgb = ORANGE
    c.line.width = Pt(1)
    txt(s, Emu(int(x) + Inches(0.12)), chip_y, Emu(int(chip_w) - Inches(0.24)), Inches(0.95),
        [[(f"{i+1}", 13, ORANGE, True)], [(ctext, 13.5, BEIGE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=Pt(2))
txt(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.5),
    [[("hanwha-atlas-production.up.railway.app", 16, BEIGE_DIM, False)]],
    align=PP_ALIGN.CENTER)
brand_tag(s)
page_num(s, 6, TOTAL)

# ---- 저장 ----
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "한화ATLAS_발표.pptx")
prs.save(out)
print("SAVED:", out)
print("slides:", len(prs.slides._sldIdLst))
